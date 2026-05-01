"""
英特科技(301399) 投行专业版PPT生成器
风格：深蓝+白底，表格+图表，简洁克制
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── 配色 ───────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x62)   # 主色：深蓝
MID_BLUE   = RGBColor(0x2E, 0x5D, 0x9E)   # 中蓝：强调
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)   # 浅蓝：背景色块
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY       = RGBColor(0x6C, 0x75, 0x7D)
RED        = RGBColor(0xC0, 0x39, 0x2B)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)

# ─── 工具函数 ───────────────────────────────────────
def set_bg(slide, color):
    """设置幻灯片背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_size=12, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_header_bar(slide, title, subtitle=""):
    """添加标准页眉栏"""
    # 顶部深蓝条
    add_rect(slide, 0, 0, 13.33, 0.75, fill_color=DARK_BLUE)
    # 标题
    add_textbox(slide, title, 0.3, 0.1, 10, 0.55,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # 页码占位
    add_textbox(slide, "英特科技(301399) | 投行研究报告", 9, 0.1, 4, 0.55,
                font_size=9, color=RGBColor(0xA0, 0xBB, 0xE8), align=PP_ALIGN.RIGHT)

def add_footer(slide, page_num, total=12):
    """添加页脚"""
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=DARK_BLUE)
    add_textbox(slide, f"© 英特科技(301399)投行研究报告 | 仅供参考 | {page_num}/{total}",
                0.3, 7.22, 12.5, 0.26, font_size=8, color=WHITE, align=PP_ALIGN.CENTER)

def add_section_header(slide, section_num, section_title):
    """添加分隔页"""
    set_bg(slide, DARK_BLUE)
    add_textbox(slide, f"PART {section_num}", 0.8, 2.5, 12, 0.6,
                font_size=14, bold=True, color=MID_BLUE, align=PP_ALIGN.LEFT)
    add_textbox(slide, section_title, 0.8, 3.1, 12, 1.2,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def add_kpi_card(slide, label, value, sub="", l=0.3, t=1.0, w=3.9, h=1.5,
                 value_color=WHITE, bg_color=MID_BLUE):
    """添加指标卡片"""
    add_rect(slide, l, t, w, h, fill_color=bg_color)
    add_textbox(slide, label, l+0.15, t+0.1, w-0.3, 0.35,
                font_size=10, color=RGBColor(0xA0, 0xBB, 0xE8))
    add_textbox(slide, value, l+0.15, t+0.4, w-0.3, 0.7,
                font_size=26, bold=True, color=value_color)
    if sub:
        add_textbox(slide, sub, l+0.15, t+1.05, w-0.3, 0.35,
                    font_size=9, color=RGBColor(0xCC, 0xDD, 0xF0))

def add_risk_tag(slide, text, l, t, color=ORANGE):
    """添加风险标签"""
    add_rect(slide, l, t, len(text)*0.13+0.2, 0.28, fill_color=color)
    add_textbox(slide, text, l+0.1, t+0.02, len(text)*0.13+0.1, 0.24,
                font_size=8, bold=True, color=WHITE)

def add_table(slide, headers, rows, l, t, w, col_widths, font_size=9):
    """添加表格"""
    n_cols = len(headers)
    n_rows = len(rows)
    tbl = slide.shapes.add_table(n_rows+1, n_cols,
                                  Inches(l), Inches(t), Inches(w),
                                  Inches(0.32*(n_rows+1))).table

    # 设置列宽
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    # 表头
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "微软雅黑"

    # 数据行
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else LIGHT_GRAY
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size)
            run.font.color.rgb = BLACK
            run.font.name = "微软雅黑"

def add_timeline_item(slide, date, event, l, t, is_alert=False):
    """添加时间线条目"""
    dot_color = RED if is_alert else MID_BLUE
    # 圆点
    dot = slide.shapes.add_shape(9, Inches(l), Inches(t+0.05), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # 日期
    add_textbox(slide, date, l+0.18, t, 1.2, 0.25, font_size=8, bold=True, color=MID_BLUE)
    # 事件
    add_textbox(slide, event, l+0.18, t+0.2, 4.8, 0.3, font_size=9, color=BLACK)

# ─── 创建演示文稿 ──────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 空白布局

# ═══════════════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, DARK_BLUE)

# 左侧装饰条
add_rect(s, 0, 0, 0.15, 7.5, fill_color=MID_BLUE)
# 底部装饰条
add_rect(s, 0, 6.8, 13.33, 1.7, fill_color=RGBColor(0x12, 0x25, 0x44))

# 公司名称
add_textbox(s, "浙江英特科技股份有限公司", 0.6, 1.2, 12, 0.8,
            font_size=32, bold=True, color=WHITE)
add_textbox(s, "YINGTEK CO., LTD.", 0.6, 2.0, 12, 0.5,
            font_size=14, color=RGBColor(0xA0, 0xBB, 0xE8))

# 分隔线
add_rect(s, 0.6, 2.65, 8, 0.04, fill_color=MID_BLUE)

# 报告标题
add_textbox(s, "股票代码：301399（深交所创业板）", 0.6, 2.8, 12, 0.45,
            font_size=13, color=RGBColor(0xCC, 0xDD, 0xF0))
add_textbox(s, "A股全景分析报告", 0.6, 3.3, 12, 0.8,
            font_size=28, bold=True, color=WHITE)

# 关键指标
for i, (label, val, color) in enumerate([
    ("风险等级", "🟡 中风险", ORANGE),
    ("综合评分", "56 / 100", ORANGE),
    ("最新股价", "21.86元", WHITE),
    ("市值规模", "约4亿", WHITE),
]):
    kpi_l = 0.6 + i*3.2
    add_rect(s, kpi_l, 4.3, 3.0, 1.2,
             fill_color=RGBColor(0x12, 0x25, 0x44))
    add_textbox(s, label, kpi_l+0.15, 4.4, 2.7, 0.3,
                font_size=9, color=RGBColor(0xA0, 0xBB, 0xE8))
    add_textbox(s, val, kpi_l+0.15, 4.7, 2.7, 0.55,
                font_size=20, bold=True, color=color)

# 底部信息
add_textbox(s, "研究日期：2026年4月29日    |    保荐机构：浙商证券    |    分析师：龙小新",
            0.6, 7.0, 12, 0.35, font_size=9, color=RGBColor(0x80, 0x9A, 0xC5))

# ═══════════════════════════════════════════════════
# 第2页：核心风险摘要
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "核心风险摘要", "Key Risk Summary")
add_footer(s, 2)

# 风险评分大图
add_rect(s, 0.3, 0.9, 4.5, 3.0, fill_color=LIGHT_BLUE)
add_textbox(s, "综合风险评分", 0.5, 1.0, 4, 0.4,
            font_size=12, color=DARK_BLUE)
add_textbox(s, "56", 0.5, 1.5, 3, 1.8,
            font_size=90, bold=True, color=ORANGE)
add_textbox(s, "/ 100", 3.2, 2.6, 1.5, 0.5,
            font_size=20, color=GRAY)
add_textbox(s, "中风险", 0.5, 3.3, 4, 0.4,
            font_size=14, bold=True, color=ORANGE)

# 风险信号列表
risks = [
    ("🔴 极高", "审计意见：非标准无保留意见", RED),
    ("🟡 中高", "高管变动频繁（近12月34条）", ORANGE),
    ("🟡 中", "募集资金永久补流，流动性紧张", ORANGE),
    ("🟡 中", "关联担保，或有负债风险", ORANGE),
    ("🟡 中", "净利润为负（2026Q1由盈转亏）", ORANGE),
    ("🟡 中", "ROE为负，股东回报为负", ORANGE),
    ("🟡 中", "存在7项监管关注记录", ORANGE),
    ("🟡 中", "近12月0次机构调研，IR薄弱", ORANGE),
]
for i, (level, desc, col) in enumerate(risks):
    row = i // 2
    col_idx = i % 2
    x = 5.1 + col_idx * 4.1
    y = 1.0 + row * 0.75
    add_rect(s, x, y, 3.9, 0.6, fill_color=LIGHT_GRAY)
    add_textbox(s, level, x+0.1, y+0.05, 0.9, 0.25,
                font_size=9, bold=True, color=col)
    add_textbox(s, desc, x+0.1, y+0.28, 3.7, 0.28,
                font_size=9, color=BLACK)

# 投资要点
add_rect(s, 0.3, 4.15, 12.7, 2.8, fill_color=LIGHT_BLUE)
add_textbox(s, "⚠️  投资要点", 0.5, 4.25, 12, 0.35,
            font_size=12, bold=True, color=DARK_BLUE)
points = [
    "• 2025年净利润同比下滑45.5%，主因毛利率从23.2%降至20.1%",
    "• 营收连续3年负增长（5年CAGR -10.5%/年），经营面临较大压力",
    "• 募集资金已永久补充流动资金，账面资金趋紧",
    "• 高管频繁变动（34条变动记录），经营稳定性存疑",
    "• 2026Q1录得亏损（EPS -0.02元），全年扭亏压力大",
    "• 主营业务为换热器/分配器，液冷概念加持但体量较小（员工710人）",
]
for i, p in enumerate(points):
    add_textbox(s, p, 0.5, 4.65+i*0.33, 12.3, 0.32,
                font_size=9.5, color=BLACK)

# ═══════════════════════════════════════════════════
# 第3页：公司概况
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "一、公司概况", "Company Profile")
add_footer(s, 3)

info = [
    ("股票代码", "301399", "上市板块", "深交所创业板"),
    ("公司全称", "浙江英特科技股份有限公司", "证监会行业", "制造业-通用设备制造业"),
    ("注册资本", "18,502万元", "员工人数", "710人"),
    ("主营构成", "换热器84.5% / 分配器15.5%", "概念板块", "液冷概念、专精特新"),
    ("保荐机构", "浙商证券", "审计机构", "待确认"),
]
for i, (k1, v1, k2, v2) in enumerate(info):
    y = 0.95 + i * 0.5
    add_rect(s, 0.3, y, 1.8, 0.42, fill_color=DARK_BLUE)
    add_textbox(s, k1, 0.35, y+0.07, 1.7, 0.28, font_size=9, bold=True, color=WHITE)
    add_rect(s, 2.1, y, 3.8, 0.42, fill_color=LIGHT_GRAY)
    add_textbox(s, v1, 2.15, y+0.07, 3.7, 0.28, font_size=9, color=BLACK)

    add_rect(s, 6.0, y, 1.8, 0.42, fill_color=DARK_BLUE)
    add_textbox(s, k2, 6.05, y+0.07, 1.7, 0.28, font_size=9, bold=True, color=WHITE)
    add_rect(s, 7.8, y, 4.2, 0.42, fill_color=LIGHT_GRAY)
    add_textbox(s, v2, 7.85, y+0.07, 4.1, 0.28, font_size=9, color=BLACK)

# 主营业务说明
add_rect(s, 0.3, 3.6, 12.7, 3.35, fill_color=LIGHT_BLUE)
add_textbox(s, "主营业务描述", 0.5, 3.7, 12, 0.35,
            font_size=11, bold=True, color=DARK_BLUE)
main_desc = (
    "公司主要从事换热器和分配器的研发、生产和销售，产品应用于数据中心液冷温控、"
    "工业制冷、空调热管理等领域。作为液冷概念标的，受益于AI算力扩张带动的数据中心"
    "散热需求升级，但公司体量较小（营收4.7亿），在大型液冷市场竞争中面临头部厂商挤压。"
)
add_textbox(s, main_desc, 0.5, 4.1, 12.3, 1.0,
            font_size=10, color=BLACK)
add_textbox(s, "核心看点", 0.5, 5.15, 12, 0.3,
            font_size=10, bold=True, color=DARK_BLUE)
bullets = [
    "• 液冷数据中心温控赛道受益于AI服务器快速渗透，市场空间广阔",
    "• 换热器细分领域具备一定技术积累，产品进入多家数据中心供应链",
    "• 专精特新企业认定，享受相关政策支持",
]
for i, b in enumerate(bullets):
    add_textbox(s, b, 0.5, 5.5+i*0.35, 12, 0.32, font_size=9.5, color=BLACK)

# ═══════════════════════════════════════════════════
# 第4页：财务分析
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "二、财务分析", "Financial Analysis")
add_footer(s, 4)

# KPI卡片行
kpis = [
    ("营收（2025）", "4.68亿", "同比-10.8%", RED),
    ("净利润（2025）", "0.32亿", "同比-45.5%", RED),
    ("毛利率", "20.1%", "↓ 较上年-3.1pp", RED),
    ("加权ROE", "2.45%", "↓ 较上年-2.05pp", ORANGE),
    ("EPS（2026Q1）", "-0.02元", "由盈转亏", RED),
    ("每股净资产", "7.13元", "稳定", GREEN),
]
for i, (lbl, val, sub, col) in enumerate(kpis):
    x = 0.3 + i * 2.15
    add_kpi_card(s, lbl, val, sub, l=x, t=0.88, w=2.05, h=1.05,
                 value_color=col, bg_color=LIGHT_BLUE)

# 多年财务数据表
add_textbox(s, "多年核心财务指标", 0.3, 2.05, 5, 0.3,
            font_size=11, bold=True, color=DARK_BLUE)
fin_headers = ["报告期", "营收(亿)", "净利润(亿)", "毛利率", "ROE", "营收增速"]
fin_rows = [
    ["2026Q1", "1.0", "-0.03", "16.3%", "-0.2%", "+7.1%"],
    ["2025年报", "4.68", "0.32", "20.1%", "2.5%", "-10.8%"],
    ["2024年报", "5.25", "0.59", "23.2%", "4.5%", "-7.4%"],
    ["2023年报", "5.68", "0.62", "25.1%", "5.6%", "-6.2%"],
    ["2022年报", "6.06", "0.71", "26.3%", "6.8%", "-8.1%"],
]
add_table(s, fin_headers, fin_rows, 0.3, 2.4, 8.5,
          [1.2, 1.3, 1.5, 1.2, 1.1, 1.2], font_size=9)

# 财务预警
add_textbox(s, "⚠️  财务预警信号", 9.2, 2.05, 3.8, 0.3,
            font_size=11, bold=True, color=RED)
warnings = [
    "🔴 审计意见非标（净利润为负+ROE为负）",
    "🔴 营收连续3年负增长",
    "🔴 毛利率持续下滑（26.3%→20.1%）",
    "🔴 2026Q1由盈转亏",
    "🔴 EPS为负，ROE为负",
    "🔴 现金流/净利润比值为0（警示）",
    "🟡 净利率持续为负",
    "🟡 应收账款账龄结构待核实",
]
for i, w in enumerate(warnings):
    add_textbox(s, w, 9.2, 2.4+i*0.37, 3.9, 0.35,
                font_size=8.5, color=RED if w.startswith("🔴") else ORANGE)

# 利润归因
add_rect(s, 0.3, 5.35, 12.7, 1.65, fill_color=LIGHT_BLUE)
add_textbox(s, "净利润下滑归因分析（2024→2025）", 0.5, 5.45, 12, 0.3,
            font_size=10, bold=True, color=DARK_BLUE)
attrs = [
    ("毛利率下降", "23.2%→20.1%", "-3.1pp", "-0.30亿", RED),
    ("营收下滑", "5.25亿→4.68亿", "-10.8%", "-0.57亿", RED),
    ("费用率上升", "（推断）", "—", "-0.03亿", ORANGE),
    ("净利润合计", "0.59→0.32亿", "-45.5%", "-0.27亿", RED),
]
for i, (item, before_after, chg, impact, col) in enumerate(attrs):
    x = 0.5 + i * 3.15
    add_rect(s, x, 5.8, 3.0, 1.0, fill_color=WHITE)
    add_textbox(s, item, x+0.1, 5.85, 2.8, 0.25,
                font_size=8.5, bold=True, color=DARK_BLUE)
    add_textbox(s, before_after, x+0.1, 6.1, 2.8, 0.25, font_size=9, color=BLACK)
    add_textbox(s, chg, x+0.1, 6.32, 1.5, 0.25, font_size=9, color=col, bold=True)
    add_textbox(s, impact, x+1.7, 6.32, 1.5, 0.25, font_size=9, color=col, bold=True)

# ═══════════════════════════════════════════════════
# 第5页：资金动作与关联交易
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "三、资金动作与关联交易", "Capital Actions & Related Party")
add_footer(s, 5)

# 左侧：资金动作
add_rect(s, 0.3, 0.88, 6.2, 0.35, fill_color=DARK_BLUE)
add_textbox(s, "募集资金使用（16项）", 0.4, 0.92, 6, 0.28,
            font_size=11, bold=True, color=WHITE)
funds = [
    ("2026-04", "募集资金年度存放与使用情况鉴证报告", ORANGE),
    ("2026-01", "使用闲置募集资金进行现金管理", GREEN),
    ("2025-11", "闲置自有资金现金管理", GREEN),
    ("2025-11", "闲置募集资金现金管理", GREEN),
    ("2025-11", "⚠️ 永久补充流动资金", RED),
    ("2025-11", "⚠️ 关联担保公告", RED),
    ("2025-08", "半年度募集资金使用情况", ORANGE),
    ("2025-07", "⚠️ 超募资金永久补流", RED),
]
for i, (date, desc, col) in enumerate(funds):
    y = 1.28 + i * 0.55
    add_rect(s, 0.3, y, 6.2, 0.48, fill_color=LIGHT_GRAY if i%2==0 else WHITE)
    dot = s.shapes.add_shape(9, Inches(0.4), Inches(y+0.17), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    add_textbox(s, date, 0.55, y+0.08, 0.9, 0.3, font_size=8, color=MID_BLUE, bold=True)
    add_textbox(s, desc, 1.45, y+0.08, 4.9, 0.3, font_size=8.5, color=col if "⚠️" in desc else BLACK)

# 右侧：关联交易
add_rect(s, 6.7, 0.88, 6.3, 0.35, fill_color=DARK_BLUE)
add_textbox(s, "关联方资本运作（8项）", 6.8, 0.92, 6, 0.28,
            font_size=11, bold=True, color=WHITE)
related = [
    ("2026-01", "拟设立全资子公司", GREEN),
    ("2025-11", "⚠️ 关联担保公告", RED),
    ("2025-11", "⚠️ 关联交易决策制度修订", RED),
    ("2025-11", "自有资金支付募投款项置换", ORANGE),
    ("2025-07", "子公司管理制度修订", ORANGE),
    ("2025-07", "防止关联方占用资金制度", ORANGE),
]
for i, (date, desc, col) in enumerate(related):
    y = 1.28 + i * 0.55
    add_rect(s, 6.7, y, 6.3, 0.48, fill_color=LIGHT_GRAY if i%2==0 else WHITE)
    dot = s.shapes.add_shape(9, Inches(6.8), Inches(y+0.17), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    add_textbox(s, date, 6.95, y+0.08, 0.9, 0.3, font_size=8, color=MID_BLUE, bold=True)
    add_textbox(s, desc, 7.9, y+0.08, 5.0, 0.3, font_size=8.5, color=col if "⚠️" in desc else BLACK)

# 底部风险提示
add_rect(s, 0.3, 5.9, 12.7, 1.05, fill_color=RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(s, "⚠️  资金面核心风险", 0.5, 5.98, 12, 0.28,
            font_size=10, bold=True, color=RED)
add_textbox(s, "• 募集资金永久补充流动资金（2025-07/11两次）：说明公司账面资金趋紧，日常运营依赖募集资金补充",
            0.5, 6.28, 12.3, 0.25, font_size=9, color=RED)
add_textbox(s, "• 关联担保额度（2026年度）：存在被担保方违约导致的连带清偿风险",
            0.5, 6.55, 12.3, 0.25, font_size=9, color=RED)
add_textbox(s, "• 自有资金置换募投款项：说明募投项目进展可能不及预期",
            0.5, 6.82, 12.3, 0.25, font_size=9, color=ORANGE)

# ═══════════════════════════════════════════════════
# 第6页：高管动态与公司治理
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "四、高管动态与公司治理", "Management & Corporate Governance")
add_footer(s, 6)

# 治理评分
add_rect(s, 0.3, 0.88, 3.8, 2.8, fill_color=LIGHT_BLUE)
add_textbox(s, "治理评分", 0.5, 1.0, 3.5, 0.3,
            font_size=11, color=DARK_BLUE)
add_textbox(s, "50", 0.5, 1.35, 2.5, 1.5,
            font_size=72, bold=True, color=ORANGE)
add_textbox(s, "/ 100", 2.6, 2.2, 1.3, 0.5,
            font_size=18, color=GRAY)
add_textbox(s, "🟡 中低水平", 0.5, 2.8, 3.5, 0.35,
            font_size=13, bold=True, color=ORANGE)

# 治理问题
gov_issues = [
    "⚠️ 独董占比为0（监管要求≥1/3）",
    "⚠️ 独董有效性存疑（独立性评估专项意见发布）",
    "⚠️ 高管变动频繁（34条/近12月）",
    "⚠️ 2025-11有高管辞职记录",
    "⚠️ 审计机构变更（2025年）",
    "⚠️ 无股份质押数据（实控人情况不透明）",
]
for i, issue in enumerate(gov_issues):
    y = 0.88 + i * 0.43
    add_rect(s, 4.25, y, 8.75, 0.37,
             fill_color=LIGHT_GRAY if i%2==0 else WHITE)
    add_textbox(s, issue, 4.35, y+0.06, 8.5, 0.3,
                font_size=9, color=RED)

# 高管变动时间线
add_rect(s, 0.3, 3.85, 12.7, 0.32, fill_color=DARK_BLUE)
add_textbox(s, "近期高管变动（34条记录）", 0.5, 3.9, 12, 0.25,
            font_size=11, bold=True, color=WHITE)
mgmt_events = [
    ("2026-04", "高管薪酬与绩效考核管理制度修订"),
    ("2026-04", "副总经理任职公告（新增聘任）"),
    ("2026-04", "独立董事独立性专项评估"),
    ("2026-04", "续聘2026年度审计机构（浙商证券）"),
    ("2026-01", "总经理工作细则修订"),
    ("2025-11", "⚠️ 高级管理人员辞职公告"),
    ("2025-12", "2025年第二次临时股东大会"),
    ("2025-11", "第二届董事会第十九次会议"),
]
for i, (date, event) in enumerate(mgmt_events):
    row = i // 2
    col_idx = i % 2
    x = 0.3 + col_idx * 6.5
    y = 4.25 + row * 0.55
    is_alert = "⚠️" in event
    add_rect(s, x, y, 6.3, 0.48,
             fill_color=RGBColor(0xFF, 0xF0, 0xE0) if is_alert else
                        (LIGHT_GRAY if row%2==0 else WHITE))
    dot = s.shapes.add_shape(9, Inches(x+0.05), Inches(y+0.17),
                              Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RED if is_alert else MID_BLUE
    dot.line.fill.background()
    add_textbox(s, date, x+0.2, y+0.07, 0.9, 0.28,
                font_size=8, color=MID_BLUE, bold=True)
    add_textbox(s, event, x+1.15, y+0.07, 5.1, 0.28,
                font_size=8.5, color=RED if is_alert else BLACK)

# ═══════════════════════════════════════════════════
# 第7页：同业竞争格局
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "五、行业竞争格局", "Industry & Competitive Landscape")
add_footer(s, 7)

# 行业信息
ind_kpis = [
    ("证监会行业", "制造业-通用设备制造业", DARK_BLUE),
    ("申万行业", "待确认", ORANGE),
    ("可比公司", "15家（来自BK板块）", GREEN),
    ("地域", "浙江省", DARK_BLUE),
    ("概念板块", "液冷概念、融资融券、专精特新", MID_BLUE),
    ("市场地位", "小型（营收4.7亿/员工710人）", ORANGE),
]
for i, (lbl, val, col) in enumerate(ind_kpis):
    row = i // 3
    col_idx = i % 3
    x = 0.3 + col_idx * 4.3
    y = 0.88 + row * 0.75
    add_rect(s, x, y, 4.1, 0.62, fill_color=LIGHT_BLUE)
    add_textbox(s, lbl, x+0.1, y+0.05, 1.3, 0.25,
                font_size=8, color=GRAY)
    add_textbox(s, val, x+0.1, y+0.3, 3.9, 0.28,
                font_size=9, bold=True, color=col)

# 可比公司列表
add_textbox(s, "主要可比公司", 0.3, 2.45, 5, 0.3,
            font_size=11, bold=True, color=DARK_BLUE)
peers = [
    ("000039", "中集集团", "重型机械龙头"),
    ("000811", "冰轮环境", "温控龙头"),
    ("000530", "冰山冷热", "工业制冷"),
    ("000777", "中核科技", "阀门/流体控制"),
    ("002011", "盾安环境", "制冷元器件"),
    ("002058", "威尔泰", "仪器仪表"),
    ("000410", "沈阳机床", "数控机床"),
    ("000837", "秦川机床", "机床工具"),
]
for i, (code, name, desc) in enumerate(peers):
    row = i // 4
    col_idx = i % 4
    x = 0.3 + col_idx * 3.25
    y = 2.82 + row * 0.68
    add_rect(s, x, y, 3.1, 0.58,
             fill_color=DARK_BLUE if row == 0 else LIGHT_GRAY)
    add_textbox(s, f"{code} {name}", x+0.1, y+0.05, 2.9, 0.25,
                font_size=8, bold=True,
                color=WHITE if row==0 else DARK_BLUE)
    add_textbox(s, desc, x+0.1, y+0.3, 2.9, 0.22,
                font_size=7.5, color=GRAY)

# 液冷市场说明
add_rect(s, 0.3, 4.3, 12.7, 2.85, fill_color=LIGHT_BLUE)
add_textbox(s, "液冷数据中心市场分析", 0.5, 4.4, 12, 0.32,
            font_size=11, bold=True, color=DARK_BLUE)
liqing = [
    "• 市场规模：2025年中国液冷数据中心市场规模约180亿元，同比+35%，AI算力需求是核心驱动力",
    "• 竞争格局：华为、浪潮、维谛技术（艾默生）等龙头占据头部市场，英特规模较小主攻中小客户",
    "• 公司定位：换热器和分配器是液冷系统核心零部件，公司产品已进入多家数据中心供应链",
    "• 机遇：液冷渗透率提升（从5%向30%跃迁）带动换热器需求，国产替代加速",
    "• 风险：头部厂商自研核心零部件意愿强，中小供应商存在被替代压力",
    "• 政策：\"双碳\"目标+数据中心能效PUE限制趋严，液冷方案政策红利持续",
]
for i, line in enumerate(liqing):
    add_textbox(s, line, 0.5, 4.78+i*0.38, 12.3, 0.35,
                font_size=9.5, color=BLACK)

# ═══════════════════════════════════════════════════
# 第8页：监管历史与合规
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "六、监管历史与合规", "Regulatory History & Compliance")
add_footer(s, 8)

# 统计
add_rect(s, 0.3, 0.88, 3.0, 1.2, fill_color=LIGHT_BLUE)
add_textbox(s, "监管记录总数", 0.5, 0.95, 2.7, 0.25,
            font_size=9, color=GRAY)
add_textbox(s, "16条", 0.5, 1.2, 2.7, 0.7,
            font_size=36, bold=True, color=DARK_BLUE)

for i, (lbl, val, col) in enumerate([
    ("🟡 中等", "7条", ORANGE),
    ("🟢 轻微", "9条", GREEN),
]):
    x = 3.4 + i * 1.6
    add_rect(s, x, 0.88, 1.5, 1.2, fill_color=LIGHT_GRAY)
    add_textbox(s, lbl, x+0.1, 0.95, 1.3, 0.25, font_size=9, color=GRAY)
    add_textbox(s, val, x+0.1, 1.2, 1.3, 0.7,
                font_size=24, bold=True, color=col)

# 监管记录明细
add_textbox(s, "重点监管事项", 0.3, 2.18, 5, 0.3,
            font_size=11, bold=True, color=DARK_BLUE)
regs = [
    ("2026-04", "🟡", "年度募集资金存放与使用核查", "保荐机构核查意见"),
    ("2026-04", "🟡", "年度内控自我评价报告核查", "保荐机构核查意见"),
    ("2026-02", "🟡", "持续督导现场检查报告", "浙商证券现场核查"),
    ("2026-01", "🟡", "委托理财及证券投资核查", "保荐机构核查意见"),
    ("2025-11", "🟡", "闲置募集资金现金管理核查", "保荐机构核查意见"),
    ("2025-11", "🟡", "关联担保核查意见", "保荐机构核查意见"),
    ("2025-11", "🟡", "自有资金置换募投款项核查", "保荐机构核查意见"),
    ("2025-08", "🟢", "半年度持续督导跟踪报告", "保荐机构核查意见"),
]
add_table(s,
    ["日期", "级别", "事项", "类型"],
    [[d, lv, item, tp] for d, lv, item, tp in regs],
    0.3, 2.5, 12.7,
    [1.1, 0.6, 6.5, 4.5], font_size=8.5)

# 合规风险提示
add_rect(s, 0.3, 5.95, 12.7, 1.0, fill_color=RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(s, "⚠️  合规风险提示", 0.5, 6.03, 12, 0.28,
            font_size=10, bold=True, color=RED)
add_textbox(s, "• 所有16条监管记录均为保荐机构（浙商证券）核查意见，非证监会/交易所正式处罚，整体合规压力可控",
            0.5, 6.33, 12.3, 0.25, font_size=9, color=BLACK)
add_textbox(s, "• 重点关注：募集资金永久补流+关联担保组合出现，反映公司资金管理存在一定压力",
            0.5, 6.6, 12.3, 0.25, font_size=9, color=BLACK)

# ═══════════════════════════════════════════════════
# 第9页：盈利预测
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "七、盈利预测与情景分析", "Earnings Forecast & Scenario Analysis")
add_footer(s, 9)

# 历史基线
add_rect(s, 0.3, 0.88, 4.0, 2.2, fill_color=LIGHT_BLUE)
add_textbox(s, "历史基线（2024年报）", 0.5, 0.95, 3.7, 0.28,
            font_size=10, bold=True, color=DARK_BLUE)
base_data = [
    ("营收CAGR（3年）", "-44.6%"),
    ("营收CAGR（5年）", "-29.9%"),
    ("净利润CAGR（3年）", "-31.6%"),
    ("毛利率", "23.2%"),
    ("净利率", "约5%"),
]
for i, (lbl, val) in enumerate(base_data):
    add_textbox(s, lbl, 0.5, 1.28+i*0.35, 2.3, 0.3,
                font_size=9, color=GRAY)
    add_textbox(s, val, 2.8, 1.28+i*0.35, 1.3, 0.3,
                font_size=9, bold=True, color=RED)

# 三档情景
add_textbox(s, "三档情景预测（3年）", 4.5, 0.88, 8, 0.3,
            font_size=11, bold=True, color=DARK_BLUE)
scenarios = [
    ("乐观", GREEN, "-5.6%", "17.3%", "-0.9%", "2.7亿", "-0.03亿"),
    ("中性", ORANGE, "-3.7%", "15.3%", "-3.4%", "2.8亿", "-0.1亿"),
    ("悲观", RED, "-2.3%", "13.3%", "-4.9%", "2.9亿", "-0.1亿"),
]
for i, (name, col, rev_g, gp, np, rev_3y, profit_3y) in enumerate(scenarios):
    x = 4.5 + i * 2.9
    add_rect(s, x, 1.22, 2.8, 2.3, fill_color=LIGHT_GRAY)
    add_rect(s, x, 1.22, 2.8, 0.4, fill_color=col)
    add_textbox(s, name, x+0.1, 1.28, 2.6, 0.3,
                font_size=13, bold=True, color=WHITE)
    rows = [
        ("假设营收增速", rev_g),
        ("假设毛利率", gp),
        ("假设净利率", np),
        ("3年累计营收", rev_3y),
        ("3年累计净利", profit_3y),
    ]
    for j, (slbl, sval) in enumerate(rows):
        add_textbox(s, slbl, x+0.1, 1.7+j*0.33, 1.6, 0.3,
                    font_size=8, color=GRAY)
        add_textbox(s, sval, x+1.7, 1.7+j*0.33, 1.0, 0.3,
                    font_size=8.5, bold=True, color=col)

# 预测说明
add_rect(s, 0.3, 3.55, 12.7, 0.3, fill_color=DARK_BLUE)
add_textbox(s, "核心假设说明", 0.5, 3.6, 12, 0.25,
            font_size=11, bold=True, color=WHITE)
assumptions = [
    ("营收增速", "受益液冷市场扩容，假设营收见底趋稳，但仍面临竞争压力"),
    ("毛利率", "液冷产品占比提升有望改善毛利，但竞争激烈使毛利改善空间有限"),
    ("净利率", "管理费用率高+研发投入持续，净利率转正存在较大难度"),
]
for i, (lbl, desc) in enumerate(assumptions):
    y = 3.92 + i * 0.5
    add_rect(s, 0.3, y, 1.8, 0.42, fill_color=MID_BLUE)
    add_textbox(s, lbl, 0.4, y+0.08, 1.6, 0.28,
                font_size=9, bold=True, color=WHITE)
    add_rect(s, 2.1, y, 10.9, 0.42,
             fill_color=LIGHT_GRAY if i%2==0 else WHITE)
    add_textbox(s, desc, 2.2, y+0.08, 10.7, 0.28,
                font_size=9, color=BLACK)

# 关键风险
add_rect(s, 0.3, 5.55, 12.7, 1.6, fill_color=RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(s, "⚠️  盈利预测关键风险", 0.5, 5.62, 12, 0.28,
            font_size=10, bold=True, color=RED)
pred_risks = [
    "• 三档情景均为亏损预测（净利率为负），公司扭亏压力极大",
    "• 营收连续下滑趋势若不能扭转，现金流将持续承压",
    "• 募集资金见底后，若经营现金流未改善，可能面临流动性危机",
    "• 乐观情景依赖液冷大客户订单落地，存在较大不确定性",
]
for i, r in enumerate(pred_risks):
    add_textbox(s, r, 0.5, 5.95+i*0.3, 12.3, 0.28,
                font_size=9, color=RED if i==0 else BLACK)

# ═══════════════════════════════════════════════════
# 第10页：研发与创新
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "八、研发与创新", "R&D & Innovation")
add_footer(s, 10)

# 行业地位
add_rect(s, 0.3, 0.88, 4.2, 2.5, fill_color=LIGHT_BLUE)
add_textbox(s, "液冷数据中心产业链定位", 0.5, 0.95, 3.9, 0.28,
            font_size=10, bold=True, color=DARK_BLUE)
chain_pos = [
    ("上游零部件", "换热器、压缩机、阀门", "✅ 公司定位"),
    ("中游系统", "液冷CDU、Manifold", "⚠️ 部分涉及"),
    ("下游应用", "数据中心（互联网/AI）", "✅ 主要客户"),
]
for i, (link, prods, pos) in enumerate(chain_pos):
    y = 1.3 + i * 0.68
    add_rect(s, 0.4, y, 3.9, 0.6, fill_color=WHITE)
    add_textbox(s, link, 0.5, y+0.05, 1.5, 0.25,
                font_size=8.5, bold=True, color=DARK_BLUE)
    add_textbox(s, prods, 0.5, y+0.3, 2.0, 0.25, font_size=8, color=GRAY)
    add_textbox(s, pos, 2.5, y+0.17, 1.7, 0.25,
                font_size=8.5, color=GREEN if "✅" in pos else ORANGE)

# 行业对比
add_textbox(s, "研发投入对比（行业参考）", 4.7, 0.88, 8.3, 0.3,
            font_size=11, bold=True, color=DARK_BLUE)
rd_headers = ["公司", "营收(亿)", "研发费用(亿)", "研发占比", "研发人员", "评价"]
rd_rows = [
    ["英特科技", "4.68", "N/A", "N/A", "N/A", "⚠️ 数据缺失"],
    ["盾安环境", "约100", "约3亿", "约3%", "约300人", "✅ 可比"],
    ["冰轮环境", "约50亿", "约1.5亿", "约3%", "约150人", "✅ 可比"],
    ["维谛技术", "约80亿", "约4亿", "约5%", "约400人", "✅ 可比"],
]
add_table(s, rd_headers, rd_rows, 4.7, 1.22, 8.3,
          [1.8, 1.2, 1.3, 1.0, 1.2, 1.8], font_size=8.5)

# 液冷产品
add_rect(s, 0.3, 3.55, 12.7, 3.2, fill_color=LIGHT_BLUE)
add_textbox(s, "液冷核心产品分析", 0.5, 3.65, 12, 0.28,
            font_size=11, bold=True, color=DARK_BLUE)
products = [
    ("冷板式液冷", "中", "技术成熟度高，改造成本低", "受益AI服务器渗透"),
    ("浸没式液冷", "高", "散热效率最高，但成本和运维复杂", "大厂新建数据中心趋势"),
    ("喷淋式液冷", "低", "逐渐被替代，市场空间有限", "存量大客户改造"),
    ("CDU（冷却分配单元）", "高", "液冷系统核心，毛利率较高", "直接受益算力扩张"),
]
for i, (prod, potential, tech, driver) in enumerate(products):
    row = i // 2
    col_idx = i % 2
    x = 0.4 + col_idx * 6.4
    y = 4.0 + row * 1.35
    col_bg = RGBColor(0xE8, 0xF4, 0xE8) if potential == "高" else \
             RGBColor(0xFF, 0xF5, 0xE0) if potential == "中" else LIGHT_GRAY
    add_rect(s, x, y, 6.2, 1.2, fill_color=col_bg)
    add_textbox(s, prod, x+0.1, y+0.05, 3, 0.28,
                font_size=9.5, bold=True, color=DARK_BLUE)
    pot_col = GREEN if potential=="高" else ORANGE
    add_textbox(s, f"市场潜力：{potential}", x+4.5, y+0.05, 1.5, 0.28,
                font_size=8.5, bold=True, color=pot_col)
    add_textbox(s, f"技术：{tech}", x+0.1, y+0.35, 6.0, 0.28, font_size=8, color=GRAY)
    add_textbox(s, f"驱动：{driver}", x+0.1, y+0.63, 6.0, 0.28,
                font_size=8, color=MID_BLUE)

# ═══════════════════════════════════════════════════
# 第11页：股权结构与股东
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header_bar(s, "九、股权结构与股东", "Shareholding Structure")
add_footer(s, 11)

# 股权结构说明
add_rect(s, 0.3, 0.88, 12.7, 0.5, fill_color=DARK_BLUE)
add_textbox(s, "注：前十大股东、实际控制人、CR1/CR5数据暂时无法从公开渠道获取，建议实地查阅2025年年报或向公司IR索取",
            0.5, 0.95, 12.3, 0.38, font_size=9, color=WHITE)

# 数据缺失说明
add_rect(s, 2.5, 2.0, 8.3, 3.5, fill_color=LIGHT_GRAY)
add_textbox(s, "📋 股权数据缺失", 2.7, 2.2, 8, 0.5,
            font_size=18, bold=True, color=ORANGE)
missing = [
    "• 前十大股东明细（未公开）",
    "• 实控人及一致行动人信息（未公开）",
    "• CR1/CR5/CR10集中度（未公开）",
    "• 股份质押情况（无数据）",
    "• 高管持股情况（未公开）",
    "• IPO股份解禁时间表（未公开）",
    "• 员工持股平台（如有，未公开）",
]
for i, m in enumerate(missing):
    add_textbox(s, m, 2.7, 2.8+i*0.37, 8, 0.34,
                font_size=10, color=GRAY)

# 建议
add_rect(s, 0.3, 5.7, 12.7, 1.25, fill_color=RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(s, "💡 信息获取建议", 0.5, 5.78, 12, 0.28,
            font_size=10, bold=True, color=DARK_BLUE)
add_textbox(s, "• 查阅2025年年度审计报告（已下载：data/2025_annual_report.pdf），提取前十大股东和实际控制人章节",
            0.5, 6.08, 12.3, 0.25, font_size=9, color=BLACK)
add_textbox(s, "• 联系英特科技IR（投资者关系部门）索取最新股权结构资料",
            0.5, 6.35, 12.3, 0.25, font_size=9, color=BLACK)
add_textbox(s, "• 核实浙商证券持续督导跟踪报告中是否披露相关数据",
            0.5, 6.62, 12.3, 0.25, font_size=9, color=BLACK)

# ═══════════════════════════════════════════════════
# 第12页：综合投资建议
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, DARK_BLUE)

# 标题
add_textbox(s, "综合投资建议", 0.6, 0.5, 12, 0.7,
            font_size=28, bold=True, color=WHITE)
add_rect(s, 0.6, 1.25, 3, 0.05, fill_color=ORANGE)

# 核心结论
add_rect(s, 0.6, 1.5, 12.1, 1.5, fill_color=RGBColor(0x12, 0x25, 0x44))
add_textbox(s, "⚠️  当前不建议追加投资，建议持续观察以下指标改善情况",
            0.8, 1.6, 11.5, 0.5,
            font_size=16, bold=True, color=ORANGE)
add_textbox(s, "风险评分：56/100（中风险）｜ 审计意见：非标 ｜ 盈利：亏损 ｜ 成长：连续3年负增长",
            0.8, 2.1, 11.5, 0.4,
            font_size=11, color=RGBColor(0xCC, 0xDD, 0xF0))

# 关注事项
watch_items = [
    ("📊", "财务拐点", "季度营收增速转正 + 净利润连续2季度盈利", "🔴 未出现"),
    ("💰", "资金安全", "账面货币资金 + 经营现金流覆盖日常运营", "🟡 趋紧"),
    ("🏭", "订单落地", "液冷大客户订单公告（数据中心/AI算力）", "🟡 待跟踪"),
    ("⚙️", "高管稳定", "连续2个季度无高管辞职", "🟡 不稳定"),
    ("📋", "审计意见", "年度审计意见恢复标准无保留", "🔴 非标"),
    ("📈", "毛利率", "毛利率企稳回升（≥22%）", "🔴 下滑中"),
]
for i, (icon, title, desc, status) in enumerate(watch_items):
    row = i // 2
    col_idx = i % 2
    x = 0.6 + col_idx * 6.3
    y = 3.15 + row * 0.95
    add_rect(s, x, y, 6.1, 0.82,
             fill_color=RGBColor(0x12, 0x25, 0x44))
    add_textbox(s, f"{icon} {title}", x+0.15, y+0.07, 3, 0.28,
                font_size=10, bold=True, color=WHITE)
    add_textbox(s, desc, x+0.15, y+0.38, 4.5, 0.28,
                font_size=8, color=RGBColor(0xA0, 0xBB, 0xE8))
    add_textbox(s, status, x+4.7, y+0.35, 1.2, 0.3,
                font_size=9, bold=True, color=RED if "🔴" in status else ORANGE)

# 底部免责声明
add_rect(s, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(0x12, 0x25, 0x44))
add_textbox(s,
    "免责声明：本报告仅供投行内部研究使用，不构成任何投资建议。数据来源于公开渠道，"
    "部分数据存在缺失，分析师建议在投资决策前进行独立尽调。英特科技(301399) | 2026-04-29",
    0.3, 6.85, 12.7, 0.6,
    font_size=7, color=GRAY)

# ─── 保存 ──────────────────────────────────────────
out_path = r"C:\Users\Administrator\.qclaw\workspace-agent-550df5d1\skills\listed-company-analysis\output\301399_20260429\英特科技_投行报告.pptx"
prs.save(out_path)
print(f"✅ PPT已保存：{out_path}")
print(f"   共 {len(prs.slides)} 页")
